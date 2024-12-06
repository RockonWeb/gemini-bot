import asyncio
import docx2txt
import openpyxl
from pptx import Presentation  # For PowerPoint files
import chardet  # For encoding detection
import pygments  # For code highlighting
from pygments.lexers import get_lexer_by_name
from pygments.formatters import HtmlFormatter
from google.api_core.exceptions import ServiceUnavailable
from gemini_pro_bot.llm import models, default_model
from google.generativeai.types.generation_types import (
    StopCandidateException,
    BlockedPromptException,
)
from telegram import InlineKeyboardButton, InlineKeyboardMarkup
from telegram import Update
from telegram.ext import (
    ContextTypes,
)
from telegram.error import NetworkError, BadRequest
from telegram.constants import ChatAction, ParseMode
from gemini_pro_bot.html_format import format_message
import PIL.Image as Image
from io import BytesIO
import docx
import logging
import sqlite3

def init_db():
    connection = sqlite3.connect("settings.db")
    cursor = connection.cursor()
    cursor.execute(
        """CREATE TABLE IF NOT EXISTS user_settings (
           user_id INTEGER PRIMARY KEY,
           model_name TEXT
        )"""
    )
    connection.commit()
    connection.close()

def user_exists(user_id):
    connection = sqlite3.connect("settings.db")
    cursor = connection.cursor()
    cursor.execute("SELECT 1 FROM user_settings WHERE user_id = ?", (user_id,))
    result = cursor.fetchone()
    connection.close()
    return result is not None


async def new_chat(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    current_model = context.chat_data.get("current_model", default_model)
    context.chat_data["chat"] = current_model.start_chat()
    await update.message.reply_text("New chat session started.")

async def start(update: Update, _: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик команды /start."""
    user = update.effective_user
    default_model_name = default_model.model_name  # Используем имя модели по умолчанию

    # Сохранение пользователя в базу данных
    save_or_update_user_model(user.id, default_model_name)

    await update.message.reply_html(
        f"Привет, {user.mention_html()}!\n\nВыберите модель или начните отправлять сообщения.",
    )


async def set_model_buttons(update: Update, _: ContextTypes.DEFAULT_TYPE) -> None:
    """Displays buttons to select the model."""

    # Calculate number of buttons per row (adjust as needed)
    buttons_per_row = 2  # Or 3, 4, etc.

    keyboard = []
    current_row = []
    for name in models:
        current_row.append(InlineKeyboardButton(name, callback_data=f"set_model:{name}"))
        if len(current_row) == buttons_per_row:
            keyboard.append(current_row)
            current_row = []  # Start a new row

    # Add the last row if it's not full
    if current_row:
        keyboard.append(current_row)


    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text("Select a model:", reply_markup=reply_markup)

def save_or_update_user_model(user_id, model_name):
    connection = sqlite3.connect("settings.db")
    cursor = connection.cursor()
    cursor.execute(
        """INSERT INTO user_settings (user_id, model_name)
           VALUES (?, ?)
           ON CONFLICT(user_id)
           DO UPDATE SET model_name = excluded.model_name""",
        (user_id, model_name),
    )
    connection.commit()
    connection.close()


async def handle_model_selection(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = update.callback_query
    await query.answer()

    model_name = query.data.split(":")[1]
    user_id = update.effective_user.id

    if model_name in models:
        context.chat_data["current_model"] = models[model_name]
        context.chat_data["chat"] = models[model_name].start_chat()

        # Сохранение новой модели в базу данных
        save_or_update_user_model(user_id, model_name)

        await query.edit_message_text(text=f"Модель изменена на {model_name}. Новая сессия начата.")
    else:
        await query.edit_message_text(text="Выбранная модель недоступна.")


def get_all_user_settings():
    connection = sqlite3.connect("settings.db")
    cursor = connection.cursor()
    cursor.execute("SELECT user_id, model_name FROM user_settings")
    rows = cursor.fetchall()
    connection.close()

    if not rows:
        return "База данных пуста."
    
    result = "Текущие настройки пользователей:\n"
    for user_id, model_name in rows:
        result += f"user {user_id} == {model_name}\n"
    return result


async def help_command(update: Update, _: ContextTypes.DEFAULT_TYPE) -> None:
    """Send a message when the command /help is issued."""
    help_text = """
Basic commands:
/start - Start the bot
/help - Get help. Shows this message

Chat commands:
/new - Start a new chat session (model will forget previously generated messages)

/set_model - Change your model(if quota's error)

Send a message to the bot to generate a response.
"""
    await update.message.reply_text(help_text)

def load_user_model(user_id):
    connection = sqlite3.connect("settings.db")
    cursor = connection.cursor()
    cursor.execute("SELECT model_name FROM user_settings WHERE user_id = ?", (user_id,))
    result = cursor.fetchone()
    connection.close()
    return result[0] if result else None


async def new_chat(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    model_name = load_user_model(user_id)  # Загрузка из базы данных
    current_model = models.get(model_name, default_model)
    context.chat_data["current_model"] = current_model
    context.chat_data["chat"] = current_model.start_chat()
    await update.message.reply_text("New chat session started.")

async def show_db_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_settings = get_all_user_settings()
    await update.message.reply_text(user_settings)



async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handles incoming text messages from users.
    Checks if a chat session exists for the user, initializes a new session if not.
    Handles model switching and preserves chat context.
    Sends the user's message to the chat session to generate a response.
    Streams the response back to the user, handling any errors.
    """
    user_id = update.effective_user.id
    # Получение текущей модели
    current_model = context.chat_data.get("current_model", default_model)

        # Если пользователя нет в базе, добавить его с моделью по умолчанию
    if not user_exists(user_id):
        save_or_update_user_model(user_id, current_model.model_name)

    # Проверка наличия активного чата
    if context.chat_data.get("chat") is None:
        # Инициализация нового чата, если его еще нет
        context.chat_data["chat"] = current_model.start_chat()
    elif context.chat_data["chat"].model != current_model:
        # Смена модели и инициализация нового чата
        context.chat_data["chat"] = current_model.start_chat()
        await update.message.reply_text(
            f"Switched to {current_model.model_name} model. Previous chat history is lost."
        )

    # Текст сообщения пользователя
    text = update.message.text

    # Отправка уведомления пользователю
    init_msg = await update.message.reply_text(
        text="Generating...", reply_to_message_id=update.message.message_id
    )
    await update.message.chat.send_action(ChatAction.TYPING)

    # Получение текущей сессии чата
    chat = context.chat_data.get("chat")

    try:
        response = await chat.send_message_async(text, stream=True)
        full_plain_message = ""  # Initialize before the loop

        async for chunk in response:  # Only one loop!
            if chunk.text:
                full_plain_message += chunk.text  # Accumulate text inside the loop
                message = format_message(full_plain_message) # Format the message

                try:  # Handle message edits
                    await init_msg.edit_text(
                        text=message, parse_mode=ParseMode.HTML, disable_web_page_preview=True
                    )
                except BadRequest as e:  # Handle Telegram errors
                    if "Message is not modified" in str(e):
                        continue  # Ignore this error and proceed
                    else:
                        raise  # Re-raise other BadRequest errors
                    
    except (StopCandidateException, BlockedPromptException, ServiceUnavailable) as e: # Combine common errors
        await init_msg.edit_text(f"An error occurred: {e}\n\nPlease contact the administrator: @RockonWeb") # More informative message
        if hasattr(chat, 'rewind') and isinstance(e, StopCandidateException):  #Rewind only for StopCandidateException
            chat.rewind()
    except Exception as e:
        print(f"Unexpected error: {e}")
        await init_msg.edit_text(f"An unexpected error occurred: {e}\n\nPlease contact the administrator: @RockonWeb") # Show error to user
        await asyncio.sleep(0.1)

async def handle_any_file(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handles any type of file sent by the user."""

    init_msg = await update.message.reply_text("Analyzing...", reply_to_message_id=update.message.message_id)
    file = None
    prompt = update.message.caption
    file_content = None  # Store the processed file content (text, bytes, etc.)
    filename = None

    if update.message.photo:
        images = update.message.photo
        unique_images = {img.file_id[:-7]: img for img in images}
        file = await list(unique_images.values())[0].get_file()
        file_content = BytesIO(await file.download_as_bytearray())
        file_content = Image.open(file_content)  # Convert to PIL Image


    elif update.message.document:  # Process documents ONLY if not audio
        try:
            file = await update.message.document.get_file()
            mime_type = update.message.document.mime_type
            filename = update.message.document.file_name
        except Exception as e:
            print(f"Unexpected error: {e}")
            await init_msg.edit_text(f"An unexpected error occurred: {e}\n\nPlease, don't use this file format") # Show error to user
            await asyncio.sleep(0.1)


        try:
            file_bytes = await file.download_as_bytearray()

            if mime_type == 'application/pdf':
                from PyPDF2 import PdfReader
                pdf_reader = PdfReader(BytesIO(file_bytes))
                file_content = ""
                for page in pdf_reader.pages:
                    file_content += page.extract_text()

            elif mime_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                try:  # Inner try-except for DOCX and DOC
                    if mime_type == 'application/msword':
                        file_content = docx2txt.process(BytesIO(file_bytes))
                    else:
                        doc = docx.Document(BytesIO(file_bytes))
                        file_content = ""
                        for paragraph in doc.paragraphs:
                            file_content += paragraph.text + "\n"
                        for table in doc.tables:  # Extract text from tables
                            for row in table.rows:
                                for cell in row.cells:
                                    file_content += cell.text + "\n"
                except Exception as e:
                    file_content = f"Error processing Word file: {e}" # Provide error message to user
            
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                workbook = openpyxl.load_workbook(BytesIO(file_bytes), read_only=True)  # Read-only for efficiency
                file_content = ""
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows():
                        for cell in row:
                            file_content += str(cell.value or "") + "\t"  # Tab-separated cells
                        file_content += "\n"

            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                presentation = Presentation(BytesIO(file_bytes))
                file_content = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_content += shape.text + "\n"

            elif filename.lower().endswith((
                '.py', '.html', '.java', '.cpp', '.js', '.c', 
                '.json', '.css', '.bat', '.toml', '.md', '.rtf',
                '.msg', '.wps', 'htm', 'odt')):
                try:
                    lexer = get_lexer_by_name(filename.split('.')[-1].lower())  # Get lexer based on file extension
                    formatter = HtmlFormatter()  # Or use another formatter if needed
                    file_content = pygments.highlight(file_bytes.decode(chardet.detect(file_bytes)['encoding']), lexer, formatter) # Detect encoding
                except Exception as e:
                    file_content = f"Could not highlight code: {e}\nRaw content:\n{file_bytes.decode(chardet.detect(file_bytes)['encoding'])}"

            elif mime_type.startswith("text/"): # Plain text files
                file_content = file_bytes.decode('utf-8') # Or other appropriate encoding

            elif mime_type.startswith("image/"): # Other image types
                file_content = BytesIO(file_bytes)
                file_content = Image.open(file_content)

            else:  # For other file types, send the raw bytes
                file_content = file_bytes  # Store as bytes

            # Combine prompt and file content if both are text
            if prompt and isinstance(file_content, str):
                prompt += f"\n\nСодержимое файла:\n{file_content}"
            elif isinstance(file_content, str) and not prompt: # Use file content if no caption
                prompt = f"Отвечай на русском. {file_content}"

        except Exception as e:
            await init_msg.edit_text(f"Error processing file: {e}\n\nPlease contact the administrator: @RockonWeb")
            logging.error(f"Error processing file: {e}")  # Log the error
            return

    if file or file_content:  # Proceed if a file or file content exists

        try:

            current_model = context.chat_data.get("current_model", default_model) # Get selected model or default

             # Note: Adapt this part to send the correct input to your model.
            if isinstance(file_content, BytesIO): # For image data
                response = await current_model.generate_content_async([prompt or "Отвечай на русском. Проанализируй", file_content], stream=True)
            elif isinstance(file_content, str): #  For other binary files
                response = await current_model.generate_content_async([prompt or "Отвечай на русском. Проанализируй", file_content], stream=True)
            else: # For text data (extracted from pdf or docx)
                response = await current_model.generate_content_async([prompt or "Отвечай на русском. Проанализируй", file_content], stream=True)

            # Отправка ответа. Проверка filename только здесь!
            if filename and filename.lower().endswith((
                '.py', '.html', '.java', '.cpp', '.js', '.c', 
                '.json', '.css', '.bat', '.toml', '.md','.rtf',
                '.msg', '.wps', 'htm', 'odt')):
                response = await current_model.generate_content_async([prompt or "Отвечай на русском. Проанализируй", file_content], stream=True)

            else:  # Отправка текстового ответа для всех остальных случаев
                response = await current_model.generate_content_async([prompt or "Отвечай на русском. Проанализируй", file_content], stream=True)

        except Exception as e:

            await init_msg.edit_text(f"The bot is currently overloaded. Please try again later.\n\n{e}")
            logging.error(f"Error during processing: {e}")
            return

    else:
        await init_msg.edit_text("Please provide a valid file (image, PDF, DOCX, TXT or other).")


    full_plain_message = ""
    async for chunk in response:
        try:
            if chunk.text:
                full_plain_message += chunk.text
                message = format_message(full_plain_message)
                init_msg = await init_msg.edit_text(
                    text=message,
                    parse_mode=ParseMode.HTML,
                    disable_web_page_preview=True,
                )
        except StopCandidateException:
            await init_msg.edit_text("The model unexpectedly stopped generating.\n\nPlease contact the administrator: @RockonWeb")
        except BadRequest:
            await response.resolve()
            continue
        except NetworkError:
            await init_msg.edit_text("Network error. Please try again later.")
            await response.resolve()
            return
        except IndexError:
            await init_msg.reply_text(
                "Some index error occurred. This response is not supported.\n\nPlease contact the administrator: @RockonWeb"
            )
            await response.resolve()
            continue
        except Exception as e:
            print(e)
            if chunk.text:
                full_plain_message = chunk.text
                message = format_message(full_plain_message)
                init_msg = await update.message.reply_text(
                    text=message,
                    parse_mode=ParseMode.HTML,
                    reply_to_message_id=init_msg.message_id,
                    disable_web_page_preview=True,
                )
        await asyncio.sleep(0.1)
